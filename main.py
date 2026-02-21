from fastapi import FastAPI
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from openai import OpenAI
import tempfile
import os

app = FastAPI()

# ✅ CORS — allow your Squarespace domain
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://www.stratumadmissions.com"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

class SongRequest(BaseModel):
    lyrics: str
    title: str
    artist: str

@app.post("/generate")
@app.post("/generate")
async def generate_ppt(data: SongRequest):

    prs = Presentation("Working Template.pptx")

    # --------- SLIDE 1 (Title Slide) ---------
    title_slide = prs.slides[0]
    title_slide.shapes.title.text = data.title

    for shape in title_slide.shapes:
        if shape.has_text_frame and shape != title_slide.shapes.title:
            shape.text_frame.text = f"By {data.artist}"
            break

    # --------- PREPARE LYRIC BLOCKS ---------
    blocks = [block.strip() for block in data.lyrics.split("\n\n") if block.strip()]

    lyric_template_slide = prs.slides[1]

    from copy import deepcopy
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.util import Pt
    from pptx.oxml.xmlchemy import OxmlElement

    def duplicate_slide(template_slide):
        new_slide = prs.slides.add_slide(template_slide.slide_layout)
        spTree = new_slide.shapes._spTree
        for el in list(spTree):
            if el.tag.endswith('}extLst'):
                continue
            spTree.remove(el)
        for el in template_slide.shapes._spTree:
            if el.tag.endswith('}extLst'):
                continue
            spTree.append(deepcopy(el))
        return new_slide

    def remove_bullets(paragraph):
        paragraph.level = 0
        pPr = paragraph._p.get_or_add_pPr()
        for child in list(pPr):
            if child.tag.endswith(('buChar','buAutoNum','buBlip','buNone')):
                pPr.remove(child)
        buNone = OxmlElement('a:buNone')
        pPr.insert(0, buNone)

    def find_lyrics_box(slide):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.strip().upper() == "LYRICS":
                return shape
        return None

    # --------- FIRST LYRIC SLIDE (USE EXISTING SLIDE 2) ---------
    first_slide = prs.slides[1]
    lyrics_box = find_lyrics_box(first_slide)

    if lyrics_box:
        tf = lyrics_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        lines = blocks[0].split("\n")

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            remove_bullets(p)
            for run in p.runs:
                run.font.size = Pt(44)

    # --------- REMAINING LYRIC SLIDES ---------
    for block in blocks[1:]:
        new_slide = duplicate_slide(lyric_template_slide)
        lyrics_box = find_lyrics_box(new_slide)

        if lyrics_box:
            tf = lyrics_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            lines = block.split("\n")

            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                remove_bullets(p)
                for run in p.runs:
                    run.font.size = Pt(44)

    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)

    return FileResponse(
        temp_file.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{data.title}.pptx"
    )
